"""Generate NEC Code Inspector Project Brief presentation."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
NAVY = RGBColor(0x1E, 0x27, 0x61)
ICE_BLUE = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF7, 0xFA)
MID_GRAY = RGBColor(0x64, 0x74, 0x8B)
DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
ACCENT_TEAL = RGBColor(0x06, 0x5A, 0x82)
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)
ACCENT_AMBER = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT_RED = RGBColor(0xEF, 0x44, 0x44)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = 13.333
SLIDE_H = 7.5


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if alpha is not None:
        shape.fill.fore_color.brightness = 0
    return shape


def add_text_box(slide, x, y, w, h, text, font_size=14, color=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_bullet_list(slide, x, y, w, h, items, font_size=14, color=DARK_TEXT, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
        pf = p._pPr
        if pf is None:
            from pptx.oxml.ns import qn
            from lxml import etree
            pf = etree.SubElement(p._p, qn('a:pPr'))
        from pptx.oxml.ns import qn
        from lxml import etree
        buChar = etree.SubElement(pf, qn('a:buChar'))
        buChar.set('char', '\u2022')
    return txBox


def add_card(slide, x, y, w, h, title, body_items, accent_color=ACCENT_TEAL, title_size=16, body_size=13):
    add_rect(slide, x, y, w, h, WHITE)
    add_rect(slide, x, y, 0.06, h, accent_color)
    add_text_box(slide, x + 0.2, y + 0.12, w - 0.3, 0.4, title, font_size=title_size, color=NAVY, bold=True)
    add_bullet_list(slide, x + 0.2, y + 0.55, w - 0.4, h - 0.7, body_items, font_size=body_size, color=DARK_TEXT, spacing=Pt(4))


# ===================== SLIDE 1: TITLE =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_rect(slide, 0, 0, SLIDE_W, 0.08, ACCENT_TEAL)
add_rect(slide, 0, SLIDE_H - 0.08, SLIDE_W, 0.08, ACCENT_TEAL)

add_text_box(slide, 1, 1.8, 11, 1.2, "NEC Code Inspector", font_size=48, color=WHITE, bold=True, align=PP_ALIGN.LEFT, font_name="Calibri")
add_text_box(slide, 1, 3.1, 11, 0.8, "AR/VR Training for the National Electrical Code", font_size=24, color=ICE_BLUE, align=PP_ALIGN.LEFT)
add_text_box(slide, 1, 4.3, 11, 0.5, "NEC/NFPA 70 \u2022 2026 Edition \u2022 zSpace Inspire 2", font_size=16, color=ICE_BLUE, align=PP_ALIGN.LEFT)

add_rect(slide, 1, 5.3, 2.5, 0.04, ACCENT_TEAL)
add_text_box(slide, 1, 5.6, 5, 0.4, "Project Brief  \u2022  April 2026", font_size=14, color=MID_GRAY, align=PP_ALIGN.LEFT)

# ===================== SLIDE 2: EXECUTIVE SUMMARY =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)
add_rect(slide, 0, 0, SLIDE_W, 0.9, NAVY)
add_text_box(slide, 0.8, 0.15, 11, 0.6, "Executive Summary", font_size=32, color=WHITE, bold=True)

add_text_box(slide, 0.8, 1.2, 11, 0.8,
    "NEC Code Inspector is an AR/VR educational app for the zSpace Inspire 2 that trains students on the National Electrical Code through hands-on 3D inspection scenarios, panel design, and NEC certification assessments.",
    font_size=16, color=DARK_TEXT)

# Key stats row
stats = [
    ("zSpace Inspire 2", "Platform", ACCENT_TEAL),
    ("3 Difficulty Tiers", "CTE \u2192 Apprentice \u2192 Expert", ACCENT_GREEN),
    ("73 Scripts", "Architecture Complete", NAVY),
    ("66 NEC Articles", "Targeting 200+ at Beta", ACCENT_AMBER),
]
for i, (big, small, color) in enumerate(stats):
    cx = 0.8 + i * 3.05
    add_rect(slide, cx, 2.4, 2.8, 1.4, WHITE)
    add_rect(slide, cx, 2.4, 2.8, 0.06, color)
    add_text_box(slide, cx + 0.15, 2.6, 2.5, 0.6, big, font_size=22, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, cx + 0.15, 3.2, 2.5, 0.4, small, font_size=12, color=MID_GRAY, align=PP_ALIGN.CENTER)

# Differentiator
add_rect(slide, 0.8, 4.2, 11.7, 0.9, NAVY)
add_text_box(slide, 1.2, 4.35, 11, 0.6,
    "Only AR/VR NEC training tool with 3D stereoscopic inspection, three difficulty tiers, and direct NEC article citation practice",
    font_size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Target users
add_text_box(slide, 0.8, 5.4, 11, 0.4, "TARGET USERS", font_size=12, color=MID_GRAY, bold=True)
add_bullet_list(slide, 0.8, 5.8, 11, 1.2, [
    "CTE high school students (Beginner mode)",
    "Trade school apprentices (Standard mode)",
    "Licensed electricians & inspectors (Expert mode)"
], font_size=14, color=DARK_TEXT)

# ===================== SLIDE 3: THREE CORE MODES =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)
add_rect(slide, 0, 0, SLIDE_W, 0.9, NAVY)
add_text_box(slide, 0.8, 0.15, 11, 0.6, "Three Core Modes", font_size=32, color=WHITE, bold=True)

add_card(slide, 0.8, 1.2, 3.7, 5.5, "1. Inspection Scenarios", [
    "Inspect 3D electrical installations",
    "Flag violations and cite NEC articles",
    "Scored feedback with article review",
    "5 scenarios planned:",
    "  Residential Service Panel",
    "  Branch Circuit Wiring (12 violations)",
    "  Grounding & Bonding (10 violations)",
    "  Commercial Installation",
    "  Outdoor/Wet Location"
], accent_color=ACCENT_TEAL)

add_card(slide, 4.8, 1.2, 3.7, 5.5, "2. Panel Design Sandbox", [
    "Design residential electrical panels",
    "Drag breakers to panel slots",
    "Route wires with correct gauge",
    "10-rule NEC compliance checker",
    "Art. 220 load calculation engine",
    "Residential 200A panel challenge",
    "12 required circuits to place",
    "Score: compliance + load accuracy"
], accent_color=ACCENT_GREEN)

add_card(slide, 8.8, 1.2, 3.7, 5.5, "3. Reference & Assessment", [
    "Searchable NEC article database",
    "10 quick reference cards",
    "Progress dashboard with history",
    "Chapter mastery tracking",
    "Certificate generation",
    "Score tracking across sessions",
    "Best-attempt leaderboard",
    "JSON persistence for LMS export"
], accent_color=ACCENT_AMBER)

# ===================== SLIDE 4: DIFFICULTY TIERS =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)
add_rect(slide, 0, 0, SLIDE_W, 0.9, NAVY)
add_text_box(slide, 0.8, 0.15, 11, 0.6, "Adaptive Difficulty System", font_size=32, color=WHITE, bold=True)

headers = ["Feature", "Beginner (CTE)", "Standard (Apprentice)", "Expert (Licensed)"]
rows = [
    ["NEC Citation Method", "Dropdown menu", "Searchable field", "Free text input"],
    ["Hints & Scaffolding", "Yes + timed hints", "No", "No"],
    ["Time Limit", "None", "None", "20 minutes"],
    ["Subtle Violations", "Hidden", "Hidden", "Visible"],
    ["False Positive Penalty", "No", "No", "Yes"],
    ["Violation Count", "5 per scenario", "10 per scenario", "All (10-12)"],
]

col_widths = [2.8, 2.7, 2.9, 2.9]
table_x = 0.8
table_y = 1.3
row_h = 0.55

# Header row
cx = table_x
for j, (header, cw) in enumerate(zip(headers, col_widths)):
    color = NAVY if j == 0 else [ACCENT_TEAL, ACCENT_GREEN, ACCENT_AMBER][j-1]
    add_rect(slide, cx, table_y, cw, row_h, color)
    add_text_box(slide, cx + 0.1, table_y + 0.1, cw - 0.2, row_h - 0.2, header, font_size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    cx += cw + 0.05

# Data rows
for i, row in enumerate(rows):
    cy = table_y + (i + 1) * (row_h + 0.05)
    cx = table_x
    bg = WHITE
    for j, (cell, cw) in enumerate(zip(row, col_widths)):
        add_rect(slide, cx, cy, cw, row_h, bg)
        a = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
        add_text_box(slide, cx + 0.1, cy + 0.1, cw - 0.2, row_h - 0.2, cell, font_size=12, color=DARK_TEXT, bold=(j == 0), align=a)
        cx += cw + 0.05

add_text_box(slide, 0.8, 5.6, 11, 0.8,
    "Students progress from guided Beginner mode through self-directed Expert mode, mirroring the journey from CTE classroom to professional licensing exam.",
    font_size=14, color=MID_GRAY)

# ===================== SLIDE 5: ARCHITECTURE =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)
add_rect(slide, 0, 0, SLIDE_W, 0.9, NAVY)
add_text_box(slide, 0.8, 0.15, 11, 0.6, "Technical Architecture", font_size=32, color=WHITE, bold=True)

# Left column: tech stack
add_text_box(slide, 0.8, 1.2, 5, 0.4, "TECH STACK", font_size=12, color=MID_GRAY, bold=True)
stack_items = [
    ("Unity 6.3 LTS", "+ zCore 6.3"),
    ("Windows 11", "zSpace Inspire 2"),
    ("6DOF Stylus", "Primary input + mouse fallback"),
    ("World Space UI", "Required for stereo 3D"),
    ("90fps Target", "Stereo rendering minimum"),
    ("JSON Data", "NEC articles + progress persistence"),
]
for i, (label, detail) in enumerate(stack_items):
    cy = 1.7 + i * 0.55
    add_rect(slide, 0.8, cy, 0.06, 0.4, ACCENT_TEAL)
    add_text_box(slide, 1.05, cy, 2.2, 0.4, label, font_size=14, color=DARK_TEXT, bold=True)
    add_text_box(slide, 3.3, cy, 3, 0.4, detail, font_size=12, color=MID_GRAY)

# Right column: module map
add_text_box(slide, 7.2, 1.2, 5, 0.4, "73 SCRIPTS ACROSS 10 MODULES", font_size=12, color=MID_GRAY, bold=True)
modules = [
    ("Core", "8", ACCENT_TEAL),
    ("PanelSandbox", "11", ACCENT_GREEN),
    ("Inspection", "8", ACCENT_TEAL),
    ("UI", "7", ACCENT_AMBER),
    ("Inputs", "7", MID_GRAY),
    ("Data", "5", NAVY),
    ("Editor", "5", MID_GRAY),
    ("Tools", "4", ACCENT_GREEN),
    ("Utils", "8", MID_GRAY),
    ("NEC", "2", NAVY),
]
for i, (name, count, color) in enumerate(modules):
    cy = 1.7 + i * 0.45
    bar_w = int(count) * 0.3
    add_rect(slide, 7.2, cy, bar_w, 0.32, color)
    add_text_box(slide, 7.3, cy + 0.02, bar_w - 0.1, 0.28, count, font_size=11, color=WHITE, bold=True)
    add_text_box(slide, 7.2 + bar_w + 0.15, cy + 0.02, 3, 0.28, name, font_size=12, color=DARK_TEXT)

# ===================== SLIDE 6: TIMELINE =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)
add_rect(slide, 0, 0, SLIDE_W, 0.9, NAVY)
add_text_box(slide, 0.8, 0.15, 11, 0.6, "Development Timeline", font_size=32, color=WHITE, bold=True)

phases = [
    ("Steps 1-3", "Weeks 1-2", "Done", "Core systems, NEC database, inspection framework", ACCENT_GREEN),
    ("Steps 4-5", "Weeks 3-4", "Done", "Branch circuits (12 violations), panel sandbox", ACCENT_GREEN),
    ("Steps 6-7", "Weeks 5-6", "Done", "Reference cards, certificates, audio, main menu", ACCENT_GREEN),
    ("Unity Assembly", "Weeks 7-8", "Next", "Scenes, prefabs, 3D assets, UI wiring", ACCENT_AMBER),
    ("Alpha", "Weeks 9-10", "Planned", "Scenarios 3-5, 200+ articles, tutorial system", ACCENT_TEAL),
    ("Beta", "Weeks 11-16", "Planned", "Full content, analytics, exam prep mode", ACCENT_TEAL),
    ("Release", "Weeks 17-20", "Planned", "Polish, accessibility, state licensing alignment", NAVY),
]

for i, (phase, timing, status, desc, color) in enumerate(phases):
    cy = 1.2 + i * 0.82
    # Status dot
    dot_color = ACCENT_GREEN if status == "Done" else (ACCENT_AMBER if status == "Next" else MID_GRAY)
    add_rect(slide, 0.8, cy + 0.12, 0.18, 0.18, dot_color)
    # Vertical line
    if i < len(phases) - 1:
        add_rect(slide, 0.86, cy + 0.35, 0.06, 0.52, RGBColor(0xE2, 0xE8, 0xF0))
    # Content
    add_text_box(slide, 1.2, cy, 2.2, 0.4, phase, font_size=16, color=DARK_TEXT, bold=True)
    add_text_box(slide, 3.5, cy, 1.5, 0.4, timing, font_size=13, color=MID_GRAY)
    add_rect(slide, 5.2, cy + 0.08, 0.8, 0.28, color)
    add_text_box(slide, 5.22, cy + 0.08, 0.76, 0.28, status, font_size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, 6.2, cy, 6.5, 0.4, desc, font_size=13, color=DARK_TEXT)

# ===================== SLIDE 7: AI INTEGRATION =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_rect(slide, 0, 0, SLIDE_W, 0.08, ACCENT_TEAL)
add_text_box(slide, 0.8, 0.4, 11, 0.6, "AI Integration Opportunity", font_size=32, color=WHITE, bold=True)

# World Labs section
add_rect(slide, 0.8, 1.3, 6, 3.5, RGBColor(0x27, 0x33, 0x72))
add_text_box(slide, 1.0, 1.45, 5.5, 0.4, "World Labs Marble", font_size=20, color=ACCENT_TEAL, bold=True)
add_text_box(slide, 1.0, 1.9, 5.5, 0.4, "AI-Generated 3D Environments", font_size=14, color=ICE_BLUE)
add_bullet_list(slide, 1.0, 2.4, 5.5, 2.2, [
    "Generate photorealistic rooms from text prompts",
    "Export GLB mesh + Gaussian splat to Unity",
    "~5 min per environment, ~$1.28 each",
    "Replace manual 3D modeling for 5+ scenarios",
    "Total estimated budget: ~$250"
], font_size=13, color=WHITE, spacing=Pt(5))

# Other AI
add_rect(slide, 7.2, 1.3, 5.3, 3.5, RGBColor(0x27, 0x33, 0x72))
add_text_box(slide, 7.4, 1.45, 5, 0.4, "Additional AI Features", font_size=20, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, 7.4, 2.0, 5, 2.6, [
    "NEC Q&A Chatbot (Claude API)",
    "  Students ask contextual code questions",
    "Adaptive Difficulty (ML)",
    "  Auto-adjust based on performance",
    "Procedural Violations",
    "  Unique sessions, no memorization",
    "Photo-to-Scenario",
    "  Teachers upload photos \u2192 AI generates scenarios"
], font_size=13, color=WHITE, spacing=Pt(3))

# Bottom cost callout
add_rect(slide, 0.8, 5.2, 11.7, 0.8, ACCENT_TEAL)
add_text_box(slide, 1.2, 5.3, 11, 0.6,
    "World Labs: ~170 generations at $1.28 each = $218 total  |  Start with $50 for prototyping",
    font_size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# ===================== SLIDE 8: VALUE PROPOSITION =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)
add_rect(slide, 0, 0, SLIDE_W, 0.9, NAVY)
add_text_box(slide, 0.8, 0.15, 11, 0.6, "Value Proposition", font_size=32, color=WHITE, bold=True)

add_card(slide, 0.8, 1.2, 3.7, 3.0, "For Students", [
    "Learn NEC through hands-on 3D inspection",
    "Practice citing specific code articles",
    "Skill tested directly on licensing exams",
    "Difficulty grows from CTE to journeyman",
    "Immediate scored feedback"
], accent_color=ACCENT_TEAL, body_size=12)

add_card(slide, 4.8, 1.2, 3.7, 3.0, "For Educators", [
    "Aligns with NEC 2026 (latest code cycle)",
    "Progress tracking for grading",
    "Certificate generation for records",
    "Difficulty tiers serve mixed classrooms",
    "Scenario assessment mirrors real inspections"
], accent_color=ACCENT_GREEN, body_size=12)

add_card(slide, 8.8, 1.2, 3.7, 3.0, "For CTE Programs", [
    "Only AR/VR NEC training tool available",
    "Runs on existing zSpace Inspire 2 hardware",
    "Covers electrical curriculum requirements",
    "Certificate tracking for accreditation",
    "Differentiated offering for the program"
], accent_color=ACCENT_AMBER, body_size=12)

# Bottom: competitive differentiation
add_rect(slide, 0.8, 4.6, 11.7, 2.2, WHITE)
add_text_box(slide, 1.0, 4.75, 11, 0.4, "COMPETITIVE DIFFERENTIATION", font_size=12, color=MID_GRAY, bold=True)
add_bullet_list(slide, 1.0, 5.2, 11, 1.5, [
    "First and only AR/VR training tool for the National Electrical Code",
    "Direct NEC article citation practice \u2014 the exact skill tested on licensing exams",
    "Three difficulty tiers serve the full pipeline from CTE student to licensed electrician",
    "NEC 2026 edition ensures always-current content"
], font_size=14, color=DARK_TEXT, spacing=Pt(5))

# ===================== SLIDE 9: CONTENT COVERAGE =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)
add_rect(slide, 0, 0, SLIDE_W, 0.9, NAVY)
add_text_box(slide, 0.8, 0.15, 11, 0.6, "Content Coverage", font_size=32, color=WHITE, bold=True)

# NEC Articles
add_rect(slide, 0.8, 1.2, 5.8, 2.8, WHITE)
add_rect(slide, 0.8, 1.2, 5.8, 0.06, ACCENT_TEAL)
add_text_box(slide, 1.0, 1.35, 3, 0.4, "66 NEC Articles", font_size=20, color=NAVY, bold=True)
add_text_box(slide, 4.2, 1.35, 2.2, 0.4, "targeting 200+", font_size=12, color=MID_GRAY)
add_bullet_list(slide, 1.0, 1.85, 5.4, 2.0, [
    "Ch 1: General Requirements (Art. 110)",
    "Ch 2: Branch Circuits, Services, Grounding (Art. 210-250)",
    "Ch 3: Wiring Methods (Art. 310, 334)",
    "Ch 4: Equipment (Art. 406, 408, 410, 422, 430)",
    "Ch 6: Special Equipment (Art. 680)"
], font_size=12, color=DARK_TEXT, spacing=Pt(3))

# Violations
add_rect(slide, 7.0, 1.2, 5.5, 2.8, WHITE)
add_rect(slide, 7.0, 1.2, 5.5, 0.06, ACCENT_GREEN)
add_text_box(slide, 7.2, 1.35, 3, 0.4, "34 Violations", font_size=20, color=NAVY, bold=True)
add_text_box(slide, 10.0, 1.35, 2.2, 0.4, "across 3 scenarios", font_size=12, color=MID_GRAY)
add_bullet_list(slide, 7.2, 1.85, 5.1, 2.0, [
    "12 Branch Circuit (GFCI, AFCI, spacing, wire gauge)",
    "10 Grounding & Bonding (electrodes, bonding, GEC)",
    "12 Panel Sandbox compliance rules",
    "Difficulty-filtered: Beginner sees 4-5, Expert sees all",
    "Each violation maps to specific NEC article"
], font_size=12, color=DARK_TEXT, spacing=Pt(3))

# Quick Reference Cards
add_rect(slide, 0.8, 4.3, 11.7, 2.5, WHITE)
add_rect(slide, 0.8, 4.3, 11.7, 0.06, ACCENT_AMBER)
add_text_box(slide, 1.0, 4.45, 5, 0.4, "10 Quick Reference Cards", font_size=18, color=NAVY, bold=True)

cards = ["GFCI Protection", "AFCI Protection", "Wire Sizing", "Branch Circuits", "Receptacle Spacing",
         "Load Calculation", "Grounding", "Panel Design", "NM Cable", "2026 NEC Changes"]
for i, card in enumerate(cards):
    row = i // 5
    col = i % 5
    cx = 1.0 + col * 2.3
    cy = 5.0 + row * 0.65
    add_rect(slide, cx, cy, 2.1, 0.5, LIGHT_GRAY)
    add_text_box(slide, cx + 0.1, cy + 0.08, 1.9, 0.35, card, font_size=11, color=DARK_TEXT, align=PP_ALIGN.CENTER)

# ===================== SLIDE 10: NEXT STEPS =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_rect(slide, 0, 0, SLIDE_W, 0.08, ACCENT_TEAL)
add_rect(slide, 0, SLIDE_H - 0.08, SLIDE_W, 0.08, ACCENT_TEAL)

add_text_box(slide, 0.8, 0.6, 11, 0.8, "Next Steps", font_size=36, color=WHITE, bold=True)

add_rect(slide, 0.8, 1.6, 5.8, 4.2, RGBColor(0x27, 0x33, 0x72))
add_text_box(slide, 1.0, 1.75, 5.4, 0.4, "Immediate (Weeks 7-8)", font_size=18, color=ACCENT_TEAL, bold=True)
add_bullet_list(slide, 1.0, 2.25, 5.4, 3.0, [
    "Run 5 editor generators in Unity to create SO assets",
    "Build Boot and MainMenu scenes",
    "Assemble 3D inspection environments",
    "Wire UI prefabs to scripted panels",
    "zSpace hardware testing and 90fps validation",
    "Prototype World Labs Marble pipeline ($50)"
], font_size=14, color=WHITE, spacing=Pt(6))

add_rect(slide, 7.0, 1.6, 5.5, 4.2, RGBColor(0x27, 0x33, 0x72))
add_text_box(slide, 7.2, 1.75, 5.1, 0.4, "Decisions Needed", font_size=18, color=ACCENT_AMBER, bold=True)
add_bullet_list(slide, 7.2, 2.25, 5.1, 3.0, [
    "Approve World Labs credit budget (~$250)",
    "3D asset source: AI-generated vs. manual modeling",
    "Alpha demo target audience and date",
    "LMS integration requirements for progress export",
    "Audio asset sourcing (SFX library or custom)",
    "QA test plan for zSpace hardware matrix"
], font_size=14, color=WHITE, spacing=Pt(6))

add_text_box(slide, 0.8, 6.2, 11.7, 0.5,
    "All script architecture is complete. The path to alpha is scene assembly + 3D content.",
    font_size=16, color=ICE_BLUE, align=PP_ALIGN.CENTER)

# Save
output_path = r"C:\Users\Jilldonnelly\Documents\nec-code-inspector\docs\NEC_Code_Inspector_Project_Brief.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
